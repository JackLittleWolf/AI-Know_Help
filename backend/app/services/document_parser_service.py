import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def parse_file(filename: str, raw: bytes, chunk_size: int = 500, chunk_overlap: int = 50) -> list[dict]:
    """Parse file bytes into text chunks. Returns list of {text, chunk_index, page}."""
    ext = Path(filename).suffix.lower().lstrip(".")
    try:
        if ext == "pdf":
            text = _parse_pdf(raw)
        elif ext in ("docx", "doc"):
            text = _parse_docx(raw)
        elif ext in ("xlsx", "xls"):
            text = _parse_xlsx(raw)
        elif ext in ("pptx", "ppt"):
            text = _parse_pptx(raw)
        else:
            text = _parse_text(raw)
    except Exception as exc:
        logger.warning("Failed to parse %s: %s", filename, exc)
        text = raw.decode("utf-8", errors="replace")

    chunks = _split_text(text, chunk_size, chunk_overlap)
    return [{"text": c, "chunk_index": i, "page": 0} for i, c in enumerate(chunks)]


def _split_text(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", "。", ".", " ", ""],
        )
        return splitter.split_text(text)
    except ImportError:
        pass

    # Fallback: manual sliding window
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - chunk_overlap
        if start >= len(text):
            break
    return chunks


def _parse_pdf(raw: bytes) -> str:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(raw))
    parts = []
    for page in reader.pages:
        t = page.extract_text() or ""
        if t.strip():
            parts.append(t)
    return "\n\n".join(parts)


def _parse_docx(raw: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_xlsx(raw: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _parse_pptx(raw: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _parse_text(raw: bytes) -> str:
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")
